import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE
import math
import os

# --- config: input file and output ---
CSV_FILE = "coaches.csv"
OUT_PPTX = "coaches.pptx"
# page/card size (Letter), adjust if you use different paper/card size
PAGE_WIDTH_IN = 2.5
PAGE_HEIGHT_IN = 2

# Position map: field -> (left_in, top_in, width_in, height_in, font_pt)
# Edit coordinates to match your template's spots.
POSITIONS = {
    "name":       (0.1875, 0.1875, 1.75, 0.25, 8),
    "type":         (0.1875, 0.4375, 1, 0.125, 6),
    "salary":      (2.0625, 0.1875, 0.25, 0.25, 12),
    }

DEFAULT_FONT = "Tecmo Bowl"
FIELD_FONTS = {
}

BLUE = RGBColor(51,255,255)
GREEN = RGBColor(100,255,51)
ORANGE = RGBColor(255,153,51)

# add near POSITIONS or at top
IMAGE_POS = (1.5, 0.75, 0.75, 0.75) # (left_in, top_in, width_in, height_in)

def add_textbox(slide, text, left, top, width, height, font_pt=12, font_name=None, color_rgb=(0,0,0)):
    """Add a textbox and set font name/size/color. font_name must be installed on the system."""
    font_name = font_name or DEFAULT_FONT
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.clear()  # start empty
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = str(text)
    font = run.font
    font.name = font_name
    # ensure the latin font gets set for PowerPoint compatibility
    try:
        run._r.rPr.rFonts.set(qn('a:latin'), font_name)
    except Exception:
        pass
    font.size = Pt(font_pt)
    r, g, b = color_rgb
    font.color.rgb = RGBColor(r, g, b)

def add_oval(slide, left, top, width, height, color):
    # add an oval (position and size in inches)
    oval_left, oval_top, oval_w, oval_h = left, top, width, height
    oval = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(oval_left),
        Inches(oval_top),
        Inches(oval_w),
        Inches(oval_h)
    )
        
    # style the oval: fill color, outline, and optional text
    oval.fill.solid()
    oval.fill.fore_color.rgb = color  
    oval.line.color.rgb = color

    oval.line.width = Pt(1)
        
    # optional: add text inside the oval
    tf = oval.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "X"
    run.font.size = Pt(12)
    run.font.name = DEFAULT_FONT
    
    return oval

def add_rectangle(slide, left, top, width, height, text='X', color=RGBColor(255,255,255)):
    # add a rectangle (position and size in inches)
    rect_left, rect_top, rect_w, rect_h = left, top, width, height
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(rect_left),
        Inches(rect_top),
        Inches(rect_w),
        Inches(rect_h)
    )
        
    # style the rectangle: fill color, outline, and optional text
    rect.line.color.rgb = RGBColor(0,0,0)
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.width = Pt(1)
        
    # optional: add text inside the rectangle
    tf = rect.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(12)
    run.font.name = DEFAULT_FONT
    
    return rect

def generate_pptx():
    df = pd.read_csv(CSV_FILE)
    prs = Presentation()
    prs.slide_width = Inches(PAGE_WIDTH_IN)
    prs.slide_height = Inches(PAGE_HEIGHT_IN)
    blank_layout = prs.slide_layouts[6]  # usually blank

    for i, row in df.iterrows():
        slide = prs.slides.add_slide(blank_layout)
        
        # calculate number of nonzero attributes
        attrs = []
        attr_titles = []
        for attr in ['recruiting','tactics','development','scouting','strategy','trading','drafting','fundraising']:
            if attr in row.index and float(row[attr]) != 0:
                attrs.append(float(row[attr]))
                attr_titles.append(attr)
        
        if len(attrs) == 0:
            continue
        
        # sort both lists by attrs descending
        attrs, attr_titles = zip(*sorted(zip(attrs, attr_titles), reverse=True))
                
        if attrs[0] == 2.0:
            color = ORANGE
        elif attrs[0] == 1.0:
            color = GREEN
        elif attrs[0] == 0.5:
            color = BLUE
        
        add_oval(slide, 0.2, 0.7, 0.5, 0.5, color)
        
        if len(attrs) > 2:
            if attrs[1] == 2.0:
                color = ORANGE
            elif attrs[1] == 1.0:
                color = GREEN
            elif attrs[1] == 0.5:
                color = BLUE
                
            add_oval(slide, 1.8, 1.0, 0.5, 0.5, color)
            
            if attrs[2] == 2.0:
                color = ORANGE
            elif attrs[2] == 1.0:
                color = GREEN
            elif attrs[2] == 0.5:
                color = BLUE
                
            add_oval(slide, 0.2, 1.3, 0.5, 0.5, color)
            
            add_rectangle(slide, 0.75, 0.7, 1.5, 0.25, attr_titles[0])
            add_rectangle(slide, 1.0, 1.0, 0.75, 0.5, attr_titles[1])
            add_rectangle(slide, 0.75, 1.55, 1.5, 0.25, attr_titles[2])
            
        else:
            add_rectangle(slide, 0.75, 0.7, 1.5, 0.5, attr_titles[0])

            if len(attrs) > 1:
                if attrs[1] == 2.0:
                    color = ORANGE
                elif attrs[1] == 1.0:
                    color = GREEN
                elif attrs[1] == 0.5:
                    color = BLUE
                
            add_oval(slide, 0.2, 1.3, 0.5, 0.5, color)
            add_rectangle(slide, 0.75, 1.3, 1.5, 0.5, attr_titles[1])
                        
        # place fields using POSITIONS map; skip fields not present
        for field, pos in POSITIONS.items():
            if field in row.index:    
                left, top, w, h, fpt = pos
                font_name = FIELD_FONTS.get(field, DEFAULT_FONT)
                add_textbox(slide, row[field], left, top, w, h, font_pt=fpt, font_name=font_name)
            else:
                if field in row.index:
                    left, top, w, h, fpt = pos
                    font_name = FIELD_FONTS.get(field, DEFAULT_FONT)
                    add_textbox(slide, row[field], left, top, w, h, font_pt=fpt, font_name=font_name)

    prs.save(OUT_PPTX)
    print("Wrote", OUT_PPTX)

if __name__ == "__main__":
    generate_pptx()