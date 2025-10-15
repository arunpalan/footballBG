import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
import math
import os

# --- config: input file and output ---
CSV_FILE = "players.csv"
OUT_PPTX = "player_cards.pptx"
# page/card size (Letter), adjust if you use different paper/card size
PAGE_WIDTH_IN = 2.5
PAGE_HEIGHT_IN = 3.5

# Position map: field -> (left_in, top_in, width_in, height_in, font_pt)
# Edit coordinates to match your template's spots.
POSITIONS = {
    "player_name": (0.222, 0.222, 1.5, 0.25, 6),
    "salary":      (2.0, 0.167, 0.25, 0.25, 12),
    "power run":   (0.25, 2.9375, 0.3125, 0.3125, 12),
    "spread":      (0.8125, 2.9375, 0.3125, 0.3125, 12),
    "west coast":  (1.375, 2.9375, 0.3125, 0.3125, 12),
    "vertical":    (1.9375, 2.9375, 0.3125, 0.3125, 12),
    "clutch":      (0.375, 1.125, 0.25, 0.25, 12),
    "dev points":  (0.375, 1.75, 0.25, 0.25, 12),
}

DEFAULT_FONT = "Tecmo Bowl"
FIELD_FONTS = {
}

# add near POSITIONS or at top
IMAGE_FIELD = "image"            # CSV column name that holds image path (relative or absolute)
IMAGE_POS = (6.0, 0.5, 1.5, 1.5) # (left_in, top_in, width_in, height_in)
IMAGE_DEFAULT = None             # set to "path/to/default.png" to use a placeholder if missing

def add_textbox(slide, text, left, top, width, height, font_pt=12, font_name=None, color_rgb=(0,0,0)):
    """Add a textbox and set font name/size/color. font_name must be installed on the system."""
    font_name = font_name or DEFAULT_FONT
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.clear()  # start empty
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = str(text)
    font = run.font
    font.name = font_name
    # ensure the latin font gets set for PowerPoint compatibility
    try:
        run._r.rPr.rFonts.set(qn('a:latin'), font_name)
    except Exception:
        pass
    font.size = Pt(font_pt)
    r, g, b = color_rgb
    font.color.rgb = RGBColor(r, g, b)

def generate_pptx():
    df = pd.read_csv(CSV_FILE)
    prs = Presentation()
    prs.slide_width = Inches(PAGE_WIDTH_IN)
    prs.slide_height = Inches(PAGE_HEIGHT_IN)
    blank_layout = prs.slide_layouts[6]  # usually blank

    for i, row in df.iterrows():
        slide = prs.slides.add_slide(blank_layout)

        # Add the image to the slide, positioned at (0,0) and filling the entire slide
        left = top = Inches(0)
        
        # choose player template background image
        if row.get("side","").lower() == "offense":
            pic = slide.shapes.add_picture('player_template_offense.jpeg', left, top, 
                                    width=prs.slide_width, height=prs.slide_height)
        else:
            pic = slide.shapes.add_picture('player_template_defense.jpeg', left, top, 
                                    width=prs.slide_width, height=prs.slide_height)

        # Move the image to the background by manipulating the XML element tree
        # This sends the picture behind all other shapes on the slide
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element) 

        # add image if present / exists
        img_path = None
        if IMAGE_FIELD in row.index and pd.notna(row[IMAGE_FIELD]) and str(row[IMAGE_FIELD]).strip():
            candidate = str(row[IMAGE_FIELD]).strip()
            if os.path.isabs(candidate) and os.path.exists(candidate):
                img_path = candidate
            else:
                # try relative to script / working dir
                rel = os.path.join(os.getcwd(), candidate)
                if os.path.exists(rel):
                    img_path = rel

        if not img_path and IMAGE_DEFAULT and os.path.exists(IMAGE_DEFAULT):
            img_path = IMAGE_DEFAULT

        if img_path:
            left, top, w, h = IMAGE_POS
            try:
                slide.shapes.add_picture(img_path, Inches(left), Inches(top), width=Inches(w), height=Inches(h))
            except Exception as e:
                print(f"[Warning] could not add image for row {i}: {e}")

        # place fields using POSITIONS map; skip fields not present
        for field, pos in POSITIONS.items():
            if field in row.index:
                if field == "player_name":
                    # Combine position with name
                    name = str(row[field])
                    position = row.get("position", "")
                    to_print = f"{position}-{name}" if position else name
                elif field == "dev points":
                    fast_dev = row.get("fast dev", "")
                    if fast_dev == 'true':
                        to_print = f"{row[field]}*"
                else:
                    to_print = row[field]
                    
                left, top, w, h, fpt = pos
                font_name = FIELD_FONTS.get(field, DEFAULT_FONT)
                add_textbox(slide, to_print, left, top, w, h, font_pt=fpt, font_name=font_name)
            else:
                if field in row.index:
                    left, top, w, h, fpt = pos
                    font_name = FIELD_FONTS.get(field, DEFAULT_FONT)
                    add_textbox(slide, row[field], left, top, w, h, font_pt=fpt, font_name=font_name)

    prs.save(OUT_PPTX)
    print("Wrote", OUT_PPTX)

if __name__ == "__main__":
    generate_pptx()