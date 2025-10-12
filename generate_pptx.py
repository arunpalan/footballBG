import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt

# --- config: input file and output ---
CSV_FILE = "players_test.csv"
OUT_PPTX = "player_cards.pptx"
# page/card size (Letter), adjust if you use different paper/card size
PAGE_WIDTH_IN = 8.5
PAGE_HEIGHT_IN = 11.0

# Position map: field -> (left_in, top_in, width_in, height_in, font_pt)
# Edit coordinates to match your template's spots.
POSITIONS = {
    "player_name": (0.5, 0.5, 7.5, 0.5, 28),
    "position":    (0.5, 1.0, 3.5, 0.3, 14),
    "salary":      (4.0, 1.0, 4.0, 0.3, 14),
    "power run":   (0.5, 1.6, 3.0, 0.3, 12),
    "spread":      (3.6, 1.6, 3.0, 0.3, 12),
    "west coast":  (0.5, 2.0, 3.0, 0.3, 12),
    "vertical":    (3.6, 2.0, 3.0, 0.3, 12),
    "clutch":      (0.5, 2.6, 3.0, 0.3, 12),
    "dev points":  (3.6, 2.6, 3.0, 0.3, 12),
}

def add_textbox(slide, text, left, top, width, height, font_pt=12):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.clear()  # start empty
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(font_pt)

def generate_pptx():
    df = pd.read_csv(CSV_FILE)
    prs = Presentation()
    prs.slide_width = Inches(PAGE_WIDTH_IN)
    prs.slide_height = Inches(PAGE_HEIGHT_IN)
    blank_layout = prs.slide_layouts[6]  # usually blank

    for i, row in df.iterrows():
        slide = prs.slides.add_slide(blank_layout)
        # If you have a background template image, uncomment and set path:
        # slide.background.fill.user_picture('/path/to/template_page_image.png')

        # place fields using POSITIONS map; skip fields not present
        for field, pos in POSITIONS.items():
            if field in row.index:
                left, top, w, h, fpt = pos
                add_textbox(slide, row[field], left, top, w, h, font_pt=fpt)
            else:
                # example mapping: if CSV field uses spaces like 'power run' use that key
                if field in row.index:
                    left, top, w, h, fpt = pos
                    add_textbox(slide, row[field], left, top, w, h, font_pt=fpt)

    prs.save(OUT_PPTX)
    print("Wrote", OUT_PPTX)

if __name__ == "__main__":
    generate_pptx()