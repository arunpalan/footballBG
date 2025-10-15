import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE
import math
import os

# --- config: input file and output ---
CSV_FILE = "sponsors.csv"
OUT_PPTX = "sponsors.pptx"
# page/card size (Letter), adjust if you use different paper/card size
PAGE_WIDTH_IN = 1
PAGE_HEIGHT_IN = 1

# Position map: field -> (left_in, top_in, width_in, height_in, font_pt)
# Edit coordinates to match your template's spots.
POSITIONS = {
    "company":       (0.05, 0.05, 0.75, 0.1875, 4),
    "contract":      (0.65, 0.525, 0.1875, 0.1875, 8),
    }

DEFAULT_FONT = "Tecmo Bowl"
FIELD_FONTS = {
}

V1 = RGBColor(115,250,121)
V2 = RGBColor(212, 251, 121)
V3 = RGBColor(255,252,121)
V4 = RGBColor(255,212,121)
V5 = RGBColor(255,126,121)

# add near POSITIONS or at top
IMAGE_POS = (0.05, 0.375, 0.5, 0.5) # (left_in, top_in, width_in, height_in)

def add_textbox(slide, text, left, top, width, height, font_pt=12, font_name=None, color_rgb=(0,0,0)):
    """Add a textbox and set font name/size/color. font_name must be installed on the system."""
    font_name = font_name or DEFAULT_FONT
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.clear()  # start empty
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = str(text)
    font = run.font
    font.name = font_name
    # ensure the latin font gets set for PowerPoint compatibility
    try:
        run._r.rPr.rFonts.set(qn('a:latin'), font_name)
    except Exception:
        pass
    font.size = Pt(font_pt)
    r, g, b = color_rgb
    font.color.rgb = RGBColor(r, g, b)
    
def add_oval(slide, left, top, width, height, value):
    # add an oval (position and size in inches)
    oval_left, oval_top, oval_w, oval_h = left, top, width, height
    oval = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(oval_left),
        Inches(oval_top),
        Inches(oval_w),
        Inches(oval_h)
    )
    
    if value == 1:
        color = V1
    elif value == 2:
        color = V2
    elif value == 3:
        color = V3
    elif value == 4:
        color = V4
    else:
        color = V5
        
    # style the oval: fill color, outline, and optional text
    oval.fill.solid()
    oval.fill.fore_color.rgb = color  
    oval.line.color.rgb = color

    oval.line.width = Pt(1)
        
    # optional: add text inside the oval
    tf = oval.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = str(value)
    run.font.size = Pt(12)
    run.font.name = DEFAULT_FONT
    
    return oval

def generate_pptx():
    df = pd.read_csv(CSV_FILE)
    prs = Presentation()
    prs.slide_width = Inches(PAGE_WIDTH_IN)
    prs.slide_height = Inches(PAGE_HEIGHT_IN)
    blank_layout = prs.slide_layouts[6]  # usually blank

    for i, row in df.iterrows():
        slide = prs.slides.add_slide(blank_layout)
        
        # add image if present / exists
        add_textbox(slide, "x", 0.55, 0.575, 0.1, 0.1, 4, DEFAULT_FONT)
        
        add_oval(slide, 0.05, 0.375, 0.5, 0.5, row.get("value",""))

        # place fields using POSITIONS map; skip fields not present
        for field, pos in POSITIONS.items():
            if field in row.index:    
                left, top, w, h, fpt = pos
                font_name = FIELD_FONTS.get(field, DEFAULT_FONT)
                add_textbox(slide, row[field], left, top, w, h, font_pt=fpt, font_name=font_name)
            else:
                if field in row.index:
                    left, top, w, h, fpt = pos
                    font_name = FIELD_FONTS.get(field, DEFAULT_FONT)
                    add_textbox(slide, row[field], left, top, w, h, font_pt=fpt, font_name=font_name)

    prs.save(OUT_PPTX)
    print("Wrote", OUT_PPTX)

if __name__ == "__main__":
    generate_pptx()