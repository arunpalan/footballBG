import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
import math
import os

# --- config: input file and output ---
CSV_FILE = "strategies.csv"
OUT_PPTX = "strategies.pptx"
# page/card size (Letter), adjust if you use different paper/card size
PAGE_WIDTH_IN = 2.5
PAGE_HEIGHT_IN = 2

# Position map: field -> (left_in, top_in, width_in, height_in, font_pt)
# Edit coordinates to match your template's spots.
POSITIONS = {
    "name":       (0.1875, 0.1875, 2, 0.25, 9),
    "bonus":      (0.75, 1.5, 0.375, 0.375, 12),
    "clutch":     (2.0, 1.5, 0.375, 0.375, 12)
    }

DEFAULT_FONT = "Tecmo Bowl"
FIELD_FONTS = {
}

BASE = RGBColor(115,250,121)
INTER = RGBColor(255,252,121)
ADV = RGBColor(255,126,121)

# add near POSITIONS or at top
IMAGE_POS = (1.5, 0.75, 0.75, 0.75) # (left_in, top_in, width_in, height_in)

def add_textbox(slide, text, left, top, width, height, font_pt=12, font_name=None, color_rgb=(0,0,0)):
    """Add a textbox and set font name/size/color. font_name must be installed on the system."""
    font_name = font_name or DEFAULT_FONT
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.clear()  # start empty
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = str(text)
    font = run.font
    font.name = font_name
    # ensure the latin font gets set for PowerPoint compatibility
    try:
        run._r.rPr.rFonts.set(qn('a:latin'), font_name)
    except Exception:
        pass
    font.size = Pt(font_pt)
    r, g, b = color_rgb
    font.color.rgb = RGBColor(r, g, b)

def generate_pptx():
    df = pd.read_csv(CSV_FILE)
    prs = Presentation()
    prs.slide_width = Inches(PAGE_WIDTH_IN)
    prs.slide_height = Inches(PAGE_HEIGHT_IN)
    blank_layout = prs.slide_layouts[6]  # usually blank

    for i, row in df.iterrows():
        slide = prs.slides.add_slide(blank_layout)

        # Change background color
        background = slide.background
        fill = background.fill
        fill.solid()
        if row.get("type","").lower() == "base":
            fill.fore_color.rgb = BASE
        elif row.get("type","").lower() == "intermediate":
            fill.fore_color.rgb = INTER
        elif row.get("type","").lower() == "advanced":
            fill.fore_color.rgb = ADV
            
        # add image if present / exists
        if row.get("arch","").lower() == "vertical":
            img_path = "vertical.png"
        elif row.get("arch","").lower() == "west coast":
            img_path = "west_coast.png"
        elif row.get("arch","").lower() == "spread":
            img_path = "spread.png"
        elif row.get("arch","").lower() == "power run":
            img_path = "power_run.png"

        if img_path:
            left, top, w, h = IMAGE_POS
            try:
                slide.shapes.add_picture(img_path, Inches(left), Inches(top), width=Inches(w), height=Inches(h))
            except Exception as e:
                print(f"[Warning] could not add image for row {i}: {e}")

        add_textbox(slide, "bonus", 0.15, 1.57, 0.5, 0.375, 7, DEFAULT_FONT)
        add_textbox(slide, "c", 1.83, 1.57, 0.1, 0.375, 7, DEFAULT_FONT)
        
        # place fields using POSITIONS map; skip fields not present
        for field, pos in POSITIONS.items():
            if field in row.index:    
                left, top, w, h, fpt = pos
                font_name = FIELD_FONTS.get(field, DEFAULT_FONT)
                add_textbox(slide, row[field], left, top, w, h, font_pt=fpt, font_name=font_name)
            else:
                if field in row.index:
                    left, top, w, h, fpt = pos
                    font_name = FIELD_FONTS.get(field, DEFAULT_FONT)
                    add_textbox(slide, row[field], left, top, w, h, font_pt=fpt, font_name=font_name)

    prs.save(OUT_PPTX)
    print("Wrote", OUT_PPTX)

if __name__ == "__main__":
    generate_pptx()