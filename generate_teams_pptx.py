import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
import math
import os

# --- config: input file and output ---
CSV_FILE = "teams.csv"
OUT_PPTX = "teams.pptx"
# page/card size (Letter), adjust if you use different paper/card size
PAGE_WIDTH_IN = 2.5
PAGE_HEIGHT_IN = 2

# Position map: field -> (left_in, top_in, width_in, height_in, font_pt)
# Edit coordinates to match your template's spots.
POSITIONS = {
    "team_name":    (0.1875, 0.1875, 2, 0.25, 9),
    "off":          (0.5, 1.5625, 0.25, 0.25, 10),
    "clutch":       (1, 1.5625, 0.25, 0.25, 10),
    "power run":   (2.0625, 0.5, 0.25, 0.25, 10),
    "spread":      (2.0625, 0.8, 0.25, 0.25, 10),
    "west coast":  (2.0625, 1.1, 0.25, 0.25, 10),
    "vertical":    (2.0625, 1.4, 0.25, 0.25, 10)
    }

DEFAULT_FONT = "Tecmo Bowl"
FIELD_FONTS = {
}

# add near POSITIONS or at top
IMAGE_POS = (0.25, 0.625, 0.75, 0.75) # (left_in, top_in, width_in, height_in)

def add_textbox(slide, text, left, top, width, height, font_pt=12, font_name=None, color_rgb=(0,0,0)):
    """Add a textbox and set font name/size/color. font_name must be installed on the system."""
    font_name = font_name or DEFAULT_FONT
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.clear()  # start empty
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = str(text)
    font = run.font
    font.name = font_name
    # ensure the latin font gets set for PowerPoint compatibility
    try:
        run._r.rPr.rFonts.set(qn('a:latin'), font_name)
    except Exception:
        pass
    font.size = Pt(font_pt)
    r, g, b = color_rgb
    font.color.rgb = RGBColor(r, g, b)

def generate_pptx():
    df = pd.read_csv(CSV_FILE)
    prs = Presentation()
    prs.slide_width = Inches(PAGE_WIDTH_IN)
    prs.slide_height = Inches(PAGE_HEIGHT_IN)
    blank_layout = prs.slide_layouts[6]  # usually blank

    for i, row in df.iterrows():
        slide = prs.slides.add_slide(blank_layout)

        # Add the image to the slide, positioned at (0,0) and filling the entire slide
        left = top = Inches(0)
        
        # choose player template background image
        if row.get("quality","") == 1:
            pic = slide.shapes.add_picture('team_q1.jpeg', left, top, 
                                    width=prs.slide_width, height=prs.slide_height)
        elif row.get("quality","") == 2:
            pic = slide.shapes.add_picture('team_q2.jpeg', left, top, 
                                    width=prs.slide_width, height=prs.slide_height)
        elif row.get("quality","") == 3:
            pic = slide.shapes.add_picture('team_q3.jpeg', left, top, 
                                    width=prs.slide_width, height=prs.slide_height)
        elif row.get("quality","") == 4:
            pic = slide.shapes.add_picture('team_q4.jpeg', left, top, 
                                    width=prs.slide_width, height=prs.slide_height)
        elif row.get("quality","") == 5:
            pic = slide.shapes.add_picture('team_q5.jpeg', left, top, 
                                    width=prs.slide_width, height=prs.slide_height)
        
            
        # Move the image to the background by manipulating the XML element tree
        # This sends the picture behind all other shapes on the slide
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element) 

        # add image if present / exists
        if row.get("arch","").lower() == "vertical":
            img_path = "vertical.png"
        elif row.get("arch","").lower() == "west coast":
            img_path = "west_coast.png"
        elif row.get("arch","").lower() == "spread":
            img_path = "spread.png"
        elif row.get("arch","").lower() == "power run":
            img_path = "power_run.png"

        if img_path:
            left, top, w, h = IMAGE_POS
            try:
                slide.shapes.add_picture(img_path, Inches(left), Inches(top), width=Inches(w), height=Inches(h))
            except Exception as e:
                print(f"[Warning] could not add image for row {i}: {e}")

        # place fields using POSITIONS map; skip fields not present
        for field, pos in POSITIONS.items():
            if field in row.index:    
                left, top, w, h, fpt = pos
                font_name = FIELD_FONTS.get(field, DEFAULT_FONT)
                add_textbox(slide, row[field], left, top, w, h, font_pt=fpt, font_name=font_name)
            else:
                if field in row.index:
                    left, top, w, h, fpt = pos
                    font_name = FIELD_FONTS.get(field, DEFAULT_FONT)
                    add_textbox(slide, row[field], left, top, w, h, font_pt=fpt, font_name=font_name)

    prs.save(OUT_PPTX)
    print("Wrote", OUT_PPTX)

if __name__ == "__main__":
    generate_pptx()